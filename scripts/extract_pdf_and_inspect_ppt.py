"""Extract PDF text and inspect sample PPT structure for India AI Buildathon."""
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
PPT_DIR = ROOT / "PPT"
PDF_PATH = PPT_DIR / "Presentation Format_ India AI Impact Buildathon.pdf"
SAMPLE_PPT = PPT_DIR / "India AI Impact Buildathon - Sample PPT.pptx"
OUT_TXT = PPT_DIR / "presentation_format_extracted.txt"


def extract_pdf():
    from pypdf import PdfReader
    if not PDF_PATH.exists():
        print("PDF not found:", PDF_PATH)
        return ""
    reader = PdfReader(str(PDF_PATH))
    lines = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text:
            lines.append(f"--- Page {i+1} ---\n{text}")
    out = "\n".join(lines)
    OUT_TXT.write_text(out, encoding="utf-8")
    print("PDF extracted to", OUT_TXT)
    return out


def inspect_ppt():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    if not SAMPLE_PPT.exists():
        print("Sample PPT not found:", SAMPLE_PPT)
        return
    prs = Presentation(str(SAMPLE_PPT))
    info = ["Slides: %d" % len(prs.slides)]
    for i, slide in enumerate(prs.slides):
        info.append("\n--- Slide %d ---" % (i + 1))
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                info.append("  [%s] %s" % (shape.shape_type, shape.text.strip()[:80]))
    result = "\n".join(info)
    (PPT_DIR / "sample_ppt_structure.txt").write_text(result, encoding="utf-8")
    print("Template structure written to sample_ppt_structure.txt")
    print(result)


if __name__ == "__main__":
    extract_pdf()
    inspect_ppt()
