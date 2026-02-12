"""
Create a fixed layout PPT for India AI Impact Buildathon.
Fixes overlap issues by deleting old placeholders and creating new, precisely positioned text boxes.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent.parent
PPT_DIR = ROOT / "PPT"
TEMPLATE_PATH = PPT_DIR / "India AI Impact Buildathon - Sample PPT.pptx"
OUTPUT_PATH = PPT_DIR / "ScamShield_AI_Final_Submission.pptx"

# Slide Dimensions (10 x 5.625 inches)
WIDTH = Inches(10)
HEIGHT = Inches(5.625)

def clean_slide_content(slide):
    """
    Removes all shapes except the title and footer/slide number.
    Returns the title shape if found.
    """
    title_shape = None
    shapes_to_delete = []
    
    # Iterate through all shapes
    for shape in slide.shapes:
        # Check if it's a title
        if hasattr(shape, "text_frame") and shape == slide.shapes.title:
            title_shape = shape
            continue
        
        # Check for placeholder titles if not identified above
        try:
            if hasattr(shape, "placeholder_format"):
                if shape.placeholder_format.type == 1: # Title
                    title_shape = shape
                    continue
        except ValueError:
            pass # Not a placeholder
        
        # Keep background images or specific logos if needed (heuristic: keep images at top left)
        # But mostly we want to clear text boxes.
        if shape.has_text_frame:
             # Heuristic: if text frame contains "INDIA AI" (footer) keep it? 
             # For now, let's delete all text boxes that are not title.
             # We might lose the logo placeholder, so we'll re-add a logo placeholder if needed.
             shapes_to_delete.append(shape)
        
        # Also remove picture placeholders that are in the middle of the slide
        if shape.shape_type == 13: # PICTURE
            # If it's in the middle, likely a content placeholder
            if shape.top > Inches(1.5) and shape.top < Inches(4.5):
                shapes_to_delete.append(shape)

    # Delete shapes (iterate backwards/remove from xml to be safe, but python-pptx doesn't support del shape easy)
    # Workaround: Move them off slide or clear text. Deleting is hard in python-pptx.
    # We will Clear the text and make them invisible.
    for shape in shapes_to_delete:
        if shape.has_text_frame:
            shape.text_frame.clear()
        # Move off screen to be sure
        shape.left = Inches(20)
    
    return title_shape

def add_text_box(slide, text, left, top, width, height, font_size=12, bold=False, color=None, alignment=None):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    if bold:
        p.font.bold = True
    if color:
        p.font.color.rgb = color
    if alignment:
        p.alignment = alignment
    
    return textbox

def add_bullet_text(slide, title, bullets, left, top, width, height):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    # Title
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(6)
    
    # Bullets
    for bullet in bullets:
        p = text_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(11)
        p.level = 0
        p.space_before = Pt(3)

def main():
    print(f"Opening template: {TEMPLATE_PATH}")
    prs = Presentation(str(TEMPLATE_PATH))
    
    # =========================================================================
    # SLIDE 1: Title
    # =========================================================================
    slide = prs.slides[0]
    clean_slide_content(slide)
    
    # Re-add Title
    # Centered Title
    add_text_box(slide, "INDIA AI IMPACT BUILDATHON", Inches(1), Inches(1.5), Inches(8), Inches(1), 
                 font_size=32, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, "Challenge 2: Agentic Honeypot for Scam Detection", Inches(1), Inches(2.5), Inches(8), Inches(1), 
                 font_size=20, alignment=PP_ALIGN.CENTER, color=RGBColor(100, 100, 100))
    add_text_box(slide, "Team: ScamShield AI", Inches(1), Inches(3.5), Inches(8), Inches(0.5), 
                 font_size=16, alignment=PP_ALIGN.CENTER)

    # =========================================================================
    # SLIDE 2: Intro
    # =========================================================================
    slide = prs.slides[1]
    clean_slide_content(slide)
    
    # Logo Placeholder
    add_text_box(slide, "[PLACE TEAM LOGO HERE]", Inches(4), Inches(1), Inches(2), Inches(1), 
                 font_size=10, alignment=PP_ALIGN.CENTER)
    
    # Project Name
    add_text_box(slide, "ScamShield AI", Inches(1), Inches(2.2), Inches(8), Inches(1), 
                 font_size=36, bold=True, alignment=PP_ALIGN.CENTER, color=RGBColor(0, 51, 102))
    
    # Tagline
    tagline = ("An autonomous agentic honeypot system that detects scam messages, "
               "engages scammers with believable AI personas, and extracts actionable "
               "intelligence (UPI, Bank Accounts, Links).")
    add_text_box(slide, tagline, Inches(1.5), Inches(3.2), Inches(7), Inches(1.5), 
                 font_size=16, alignment=PP_ALIGN.CENTER)

    # =========================================================================
    # SLIDE 3: The Problem (3 Columns)
    # =========================================================================
    slide = prs.slides[2]
    clean_slide_content(slide)
    add_text_box(slide, "1) THE PROBLEM", Inches(0.5), Inches(0.4), Inches(9), Inches(0.6), font_size=24, bold=True)
    
    y_start = Inches(1.5)
    w_col = Inches(2.8)
    gap = Inches(0.3)
    h_col = Inches(3.5)
    
    # Col 1: What
    add_bullet_text(slide, "🚨 What is happening?", [
        "500,000+ scam calls/msgs daily",
        "₹60+ crore daily losses",
        "47% Indians affected",
        "Target: UPI, Loans, KYC"
    ], Inches(0.5), y_start, w_col, h_col)
    
    # Col 2: Why
    add_bullet_text(slide, "⚠️ Why is it a problem?", [
        "Passive detection fails",
        "Scammers evolve fast",
        "No intelligence gathered",
        "Citizens feel helpless"
    ], Inches(0.5 + w_col + gap), y_start, w_col, h_col)
    
    # Col 3: Who
    add_bullet_text(slide, "👥 Who is affected?", [
        "Elderly & Non-tech savvy",
        "Middle-class families",
        "Banks (Reputation loss)",
        "Law Enforcement (Overload)"
    ], Inches(0.5 + 2*(w_col + gap)), y_start, w_col, h_col)

    # =========================================================================
    # SLIDE 4: Our Solution (2x2 Grid)
    # =========================================================================
    slide = prs.slides[3]
    clean_slide_content(slide)
    add_text_box(slide, "2) OUR SOLUTION", Inches(0.5), Inches(0.4), Inches(9), Inches(0.6), font_size=24, bold=True)
    
    # Grid config
    left_m = Inches(0.8)
    top_m = Inches(1.2)
    box_w = Inches(4.0)
    box_h = Inches(1.8)
    h_gap = Inches(0.4)
    v_gap = Inches(0.3)
    
    # Box 1
    add_bullet_text(slide, "🎯 Detection", [
        "IndicBERT + Rules",
        "Hindi, English, Hinglish",
        "90%+ Accuracy Target"
    ], left_m, top_m, box_w, box_h)
    
    # Box 2
    add_bullet_text(slide, "🤖 Engagement", [
        "LangGraph Agent",
        "Dynamic Personas (Elderly, Eager)",
        "Up to 20 turns conversation"
    ], left_m + box_w + h_gap, top_m, box_w, box_h)
    
    # Box 3
    add_bullet_text(slide, "🔍 Extraction", [
        "Extracts UPI, Bank, IFSC",
        "Captures Phishing Links",
        "Structured JSON Output"
    ], left_m, top_m + box_h + v_gap, box_w, box_h)
    
    # Box 4
    add_bullet_text(slide, "🚀 Integration", [
        "REST API Endpoint",
        "Mock Scammer API Ready",
        "Dockerized & Scalable"
    ], left_m + box_w + h_gap, top_m + box_h + v_gap, box_w, box_h)

    # =========================================================================
    # SLIDE 5: How It Works (Flow)
    # =========================================================================
    slide = prs.slides[4]
    clean_slide_content(slide)
    add_text_box(slide, "3) HOW IT WORKS", Inches(0.5), Inches(0.4), Inches(9), Inches(0.6), font_size=24, bold=True)
    
    # Simple Text Flow
    flow_text = (
        "STEP 1: INPUT\n"
        "Message arrives: \"You won 10 lakh! Send OTP.\"\n\n"
        "STEP 2: DECISION (AI)\n"
        "Confidence > 70%? Yes -> Trigger Honeypot\n\n"
        "STEP 3: ENGAGE\n"
        "Agent (Persona): \"Oh wow! How do I get it?\"\n\n"
        "STEP 4: EXTRACT & OUTPUT\n"
        "Scammer shares UPI -> We capture it -> JSON Report"
    )
    add_text_box(slide, flow_text, Inches(1), Inches(1.5), Inches(8), Inches(3.5), font_size=14)
    
    # Visual cues text
    add_text_box(slide, "Simple Flow: Message In ➡️ Scam Check ➡️ Fake Victim Talks ➡️ Proof Extracted", 
                 Inches(0.5), Inches(5.0), Inches(9), Inches(0.5), font_size=12, bold=True, alignment=PP_ALIGN.CENTER, color=RGBColor(0, 102, 204))

    # =========================================================================
    # SLIDE 6: Proof It Works (List/Grid)
    # =========================================================================
    slide = prs.slides[5]
    clean_slide_content(slide)
    add_text_box(slide, "4) PROOF IT WORKS", Inches(0.5), Inches(0.4), Inches(9), Inches(0.6), font_size=24, bold=True)
    
    # 4 Key Proof Points
    proofs = [
        ("✅ Live API Demo", "POST /honeypot/engage works in <2s"),
        ("📊 Test Metrics", "Validated on 100+ scam samples"),
        ("🇮🇳 Real Context", "Handles Hindi/Hinglish perfectly"),
        ("🛡️ Extraction", "Success: Captured UPIs & Links")
    ]
    
    y_pos = Inches(1.4)
    for title, desc in proofs:
        add_text_box(slide, title, Inches(1), y_pos, Inches(3), Inches(0.5), font_size=16, bold=True)
        add_text_box(slide, desc, Inches(4.2), y_pos, Inches(5), Inches(0.5), font_size=14)
        y_pos += Inches(0.9) # Spacing

    # =========================================================================
    # SLIDE 7: Nuance (3 Blocks)
    # =========================================================================
    slide = prs.slides[6]
    clean_slide_content(slide)
    add_text_box(slide, "5) A NUANCE WE HANDLED", Inches(0.5), Inches(0.4), Inches(9), Inches(0.6), font_size=24, bold=True)
    
    y_start = Inches(1.5)
    box_h = Inches(3.5)
    box_w = Inches(2.8)
    
    add_bullet_text(slide, "🌐 Mixed Language", [
        "Handles Code-Mixing",
        "(Hindi + English)",
        "IndicBERT trained on Indian datasets"
    ], Inches(0.5), y_start, box_w, box_h)
    
    add_bullet_text(slide, "😇 Over-Polite Scams", [
        "Personas don't get suspicious",
        "Elderly/Confused tones",
        "Matches scammer's pace"
    ], Inches(0.5 + box_w + gap), y_start, box_w, box_h)
    
    add_bullet_text(slide, "🔄 Persistence", [
        "Redis Session State",
        "Remembers Context",
        "Engages for 10+ turns"
    ], Inches(0.5 + 2*(box_w + gap)), y_start, box_w, box_h)

    # =========================================================================
    # SLIDE 8: Trade-Offs (Left/Right)
    # =========================================================================
    slide = prs.slides[7]
    clean_slide_content(slide)
    add_text_box(slide, "6) TRADE-OFF & 7) FAILURE CASE", Inches(0.5), Inches(0.4), Inches(9), Inches(0.6), font_size=24, bold=True)
    
    # Left: Trade Off
    add_bullet_text(slide, "⚖️ Trade-Off: Depth vs Speed", [
        "Choice: Maximize engagement depth",
        "Why: More turns = More Intelligence",
        "Cost: Higher latency (1-2s)",
        "Not optimized for instant blocking"
    ], Inches(0.8), Inches(1.5), Inches(4), Inches(3.5))
    
    # Right: Failure Case
    add_bullet_text(slide, "❌ Failure Case: Short Convos", [
        "Struggle: Scammer stops after 1 msg",
        "Result: No intelligence extracted",
        "Mitigation: 'Eager' persona to bait",
        "Struggle: Novel templates (Low Conf)"
    ], Inches(5.2), Inches(1.5), Inches(4), Inches(3.5))

    # =========================================================================
    # SLIDE 9: Submission
    # =========================================================================
    slide = prs.slides[8]
    clean_slide_content(slide)
    
    # Title
    add_text_box(slide, "SUBMISSION DETAILS", Inches(1), Inches(1), Inches(8), Inches(0.8), 
                 font_size=28, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Content
    info_text = (
        "Team: ScamShield AI\n"
        "Challenge 2: Agentic Honeypot\n"
        "Tech Stack: IndicBERT, LangGraph, Groq, FastAPI\n\n"
        "Contact: missionupskillindia@hclguvi.com\n"
        "Subject: ScamShield AI PPT || India AI Impact Buildathon\n"
        "Deadline: 13th Feb 2026"
    )
    add_text_box(slide, info_text, Inches(2), Inches(2.2), Inches(6), Inches(3), 
                 font_size=16, alignment=PP_ALIGN.CENTER)

    # Save
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"SUCCESS: Fixed Layout PPT created: {OUTPUT_PATH}")

if __name__ == "__main__":
    main()
