"""Verify the final comprehensive PPT has all required content."""
from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
PPT_PATH = ROOT / "PPT" / "ScamShield_AI_Presentation_FINAL.pptx"

def main():
    prs = Presentation(str(PPT_PATH))
    print(f"Verifying: {PPT_PATH.name}")
    print(f"Total slides: {len(prs.slides)}\n")
    
    required_content = {
        1: ["INDIA AI IMPACT BUILDATHON", "Challenge 2"],
        2: ["ScamShield AI", "honeypot"],
        3: ["THE PROBLEM", "500,000", "scam"],
        4: ["OUR SOLUTION", "Detection", "Engagement", "Extraction", "API"],
        5: ["HOW IT WORKS", "Input", "Decision", "Output"],
        6: ["PROOF IT WORKS", "Demo", "Test"],
        7: ["NUANCE", "Mixed language", "Over-polite"],
        8: ["TRADE-OFF", "FAILURE CASE", "Engagement depth"],
        9: ["Submission by", "ScamShield AI", "Challenge 2", "Contact"]
    }
    
    all_good = True
    for i, slide in enumerate(prs.slides, 1):
        slide_text = " ".join([sh.text for sh in slide.shapes if hasattr(sh, "text")])
        print(f"Slide {i}: ", end="")
        
        if i in required_content:
            missing = []
            for keyword in required_content[i]:
                if keyword.lower() not in slide_text.lower():
                    missing.append(keyword)
            
            if missing:
                print(f"MISSING: {', '.join(missing)}")
                all_good = False
            else:
                print("OK - All required content present")
        else:
            print("OK")
    
    print(f"\n{'='*60}")
    if all_good:
        print("SUCCESS: All required content is present!")
    else:
        print("WARNING: Some content is missing. Review the slides.")
    print(f"{'='*60}")

if __name__ == "__main__":
    main()
