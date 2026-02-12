"""Inspect the generated PPT to verify content."""
from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
PPT_PATH = ROOT / "PPT" / "ScamShield_AI_India_AI_Impact_Buildathon.pptx"

def main():
    prs = Presentation(str(PPT_PATH))
    print(f"Total slides: {len(prs.slides)}\n")
    
    for i, slide in enumerate(prs.slides):
        print(f"{'='*60}")
        print(f"SLIDE {i+1}")
        print(f"{'='*60}")
        shapes_with_text = [sh for sh in slide.shapes if hasattr(sh, "text") and sh.text.strip()]
        if not shapes_with_text:
            print("  [No text content]")
        for j, shape in enumerate(shapes_with_text):
            text = shape.text.strip()
            if len(text) > 150:
                text = text[:150] + "..."
            print(f"  {j+1}. {text}")
        print()

if __name__ == "__main__":
    main()
