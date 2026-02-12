"""
Create India AI Impact Buildathon presentation for ScamShield AI.
Uses the official template and fills it with content from the PDF guidelines
and project documentation.
"""
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
PPT_DIR = ROOT / "PPT"
TEMPLATE_PATH = PPT_DIR / "India AI Impact Buildathon - Sample PPT.pptx"
OUTPUT_PATH = PPT_DIR / "ScamShield_AI_India_AI_Impact_Buildathon.pptx"

# Content per PDF guidelines and project docs (PRD, FRD, README, TECH_STACK)
SLIDE_CONTENT = {
    1: [
        ("INDIA AI IMPACT BUILDATHON", "INDIA AI IMPACT BUILDATHON\nChallenge 2: Agentic Honeypot for Scam Detection"),
    ],
    2: [
        ("gjhghjgjhg", "ScamShield AI"),
        ("[Your Logo Here]", "[Your Logo Here]"),
        ("Lorem Ipsum is simply dummy text.\nLorem Ipsum is simply dummy text of the printi",
         "Agentic honeypot system that detects scam messages, engages scammers with AI personas, and extracts actionable intelligence (UPI IDs, bank accounts, phishing links). Built for India AI Impact Buildathon 2026."),
    ],
    3: [
        ("<Heading>", "1) The Problem"),
        ("[ IMAGE ]", "[ IMAGE ]"),
        ("Lorem Ipsum is simply dummy text of the printing and typesetting industry. Lorem",
         "What is happening today: 500,000+ scam calls/messages daily in India; Rs 60+ crore daily losses; 47% of Indians affected or know victims. UPI fraud, fake loans, police/bank impersonation are predominant. Why it matters: Existing solutions only detect passively. ScamShield actively engages to extract intelligence while scammers remain unaware they are talking to AI."),
    ],
    4: [
        ("<Add your title here>", "2) Our Solution"),
        ("Lorem Ipsum is simply dummy text of the printing and typesetting industry. Lorem",
         "We built an autonomous AI honeypot that detects, engages, and extracts."),
        ("Caption\nSubheading", [
            "Detection: IndicBERT + keywords; English, Hindi, Hinglish.",
            "Engagement: Multi-turn (up to 20) with believable personas.",
            "Extraction: UPI IDs, bank accounts, IFSC, phone, phishing links.",
            "API: REST endpoint for Mock Scammer API integration.",
        ]),
    ],
    5: [
        ("<Add your title here>", "3) How It Works (Simple Flow)"),
        ("Lorem Ipsum is simply dummy text of the printing and typesetting industry. Lorem",
         "Input: A message (e.g. You won 10 lakh, share OTP). Decision: Is it a scam? If confidence > 0.7, hand off to the honeypot agent. Output: Agent replies in character; we extract UPI/bank/links and return structured JSON. No jargon: message in, scam check, fake victim talks back, we collect proof."),
        ("Use this space for highlighted text",
         "Message In -> Scam? -> Engage (persona) -> Extract (UPI/bank/links) -> JSON Out"),
    ],
    6: [
        ("<Add your title here>", "4) Proof It Works (15 seconds)"),
        ("[ IMAGE ]", "[ IMAGE ]"),
        ("Lorem Ipsum is simply dummy text of the printing and typesetting industry. Lorem",
         "Demo: POST /honeypot/engage with a scam message; get agent response and extracted intelligence. Test: 90%+ detection accuracy target; 85%+ extraction precision. Real example: Hindi/English scam phrases trigger engagement and structured extraction."),
        ("Add Pointer here", [
            "Live API demo",
            "Test results",
            "Extraction proof",
            "Metrics dashboard",
        ]),
        ("01", "01"),
        ("02", "02"),
        ("03", "03"),
        ("04", "04"),
    ],
    7: [
        ("<Heading>", "5) A Nuance We Handled"),
        ("Lorem Ipsum is simply dummy text of the printing and typesetting industry. Lorem",
         "Mixed language: We explicitly support Hindi, English, and Hinglish (code-mixed) in both detection and agent replies. Over-polite or urgent scam messages: Personas (elderly, eager, confused) are tuned to respond naturally so scammers do not get suspicious. Repeated attempts in one session: Session state in Redis keeps context so we can prolong engagement and extract more across turns."),
    ],
    8: [
        ("<Heading>", "6) Trade-Off & 7) Failure Case"),
        ("Lorem Ipsum is simply dummy text of the printing and typesetting industry. Lorem",
         "Trade-off: We chose engagement depth over containment (up to 20 turns) to maximize intelligence extraction; we accept slightly higher latency for richer responses. Failure case: Very short conversations may yield no extraction; novel scam templates not seen in training can get lower confidence. We handle these with clear confidence scores and safe fallbacks."),
        ("Description for\nSample Text", ["Trade-off", "Failure case"]),
        ("Sample Text", ["Engagement depth vs containment", "Short convos; novel scam templates"]),
    ],
    9: [
        ("Submission by:", "Submission by:"),
        ("Sample text 1\nSample text 2\nSample text 3",
         "ScamShield AI Team\nIndia AI Impact Buildathon 2026\nChallenge 2: Agentic Honeypot\nContact: missionupskillindia@hclguvi.com\nSubject: ScamShield AI PPT || India AI Impact Buildathon"),
    ],
}


def get_shape_text(shape):
    if not hasattr(shape, "text_frame"):
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs)


def replace_in_slide(slide, replacements):
    """Replace text in slide shapes. new can be str or list of str (assign in order)."""
    list_indices = {}
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        current = get_shape_text(shape)
        if not current.strip():
            continue
        for old, new in replacements:
            if old not in current:
                continue
            if isinstance(new, list):
                idx = list_indices.get(old, 0)
                if idx >= len(new):
                    continue
                replacement = new[idx]
                list_indices[old] = idx + 1
            else:
                replacement = new
            try:
                shape.text_frame.clear()
                shape.text_frame.paragraphs[0].text = replacement
            except Exception:
                pass
            break


def main():
    from pptx import Presentation

    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError("Template not found: %s" % TEMPLATE_PATH)

    prs = Presentation(str(TEMPLATE_PATH))
    if len(prs.slides) < 9:
        raise ValueError("Template has %d slides, expected 9" % len(prs.slides))

    for slide_idx, replacements in SLIDE_CONTENT.items():
        slide = prs.slides[slide_idx - 1]
        replace_in_slide(slide, replacements)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print("Saved: %s" % OUTPUT_PATH)


if __name__ == "__main__":
    main()
